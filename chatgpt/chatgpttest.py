# 导入库
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# 创建新的 PowerPoint 文件
prs = Presentation()

# 在文件中添加一张幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 在幻灯片中添加标题
# slide.shapes.title.text = '''2022 年 Java 软件工程师年终总结'''

# 在幻灯片中添加文本框
left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# 在文本框中添加文本
tf.text = "在过去的一年中，我在项目中负责了以下工作:"

# 在文本框中添加段落
p = tf.add_paragraph()
p.text = "1. 完成了项目 X 的开发，包括设计数据库架构、编写业务逻辑和前端界面。"
p = tf.add_paragraph()
p.text = "2. 参与了项目 Y 的维护，包括修复线上问题和优化系统性能。"
p = tf.add_paragraph()
p.text = "3. 协助进行了项目 Z 的代码审查，帮助团队提高代码质量。"

#在幻灯片中添加图片
left = top = Inches(1)
# pic = slide.shapes.add_picture("image.jpg", left, top)

#在幻灯片中添加图表
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'North', 'South']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7, 13.0))
chart_data.add_series('Series 2', (22.3, 28.6, 15.2, 10.2))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

#保存 PowerPoint 文件
prs.save("annual_review.pptx")