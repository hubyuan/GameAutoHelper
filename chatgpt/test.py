from pptx import Presentation

# 创建一个新的 PowerPoint 演示文稾
prs = Presentation()

# 添加一个新的幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[0])

# 在幻灯片中添加一个文本框
title = slide.shapes.title
title.text = "年终总结"

# 添加一个段落
body = slide.shapes.placeholder
tf = body.text_frame

# 在文本框中添加一些文本
tf.text = "这一年我作为一名 Java 软件工程师，取得了以下成就："

# 在文本框中添加一些段落
p = tf.add_paragraph()
p.text = "1. 完成了 XX 项目的开发"
p.level = 1

p = tf.add_paragraph()
p.text = "2. 参与了 XX 项目的设计"
p.level = 1

# 保存 PowerPoint 演示文稾
prs.save("year-end-summary.pptx")