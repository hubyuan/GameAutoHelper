from pptx import Presentation
from pptx.util import Inches

# 创建一个新的PPT文件
prs = Presentation()

# 向PPT中添加一个幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[0])

# 向幻灯片中添加文本框
left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# 在文本框中添加文本
tf.text = "Java软件工程师年终总结"

# 保存PPT文件
prs.save("annual_summary.pptx")